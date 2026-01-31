#!/usr/bin/env python3
"""
Markdown Conversion Service for Idea Database
Converts various file formats to markdown for LLM consumption
"""

import io
import os
import re
import base64
import tempfile
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
import structlog
import mimetypes

# PDF Processing
import fitz  # PyMuPDF

# Office Document Processing
from docx import Document
from pptx import Presentation

# Image Processing and OCR
from PIL import Image, ImageFilter, ImageEnhance
import pytesseract
import cv2
import numpy as np

# HTML Processing
from markdownify import markdownify as md
from bs4 import BeautifulSoup
import chardet

logger = structlog.get_logger()

class MarkdownConverter:
    """Service for converting various file formats to markdown"""
    
    def __init__(self, config: dict = None):
        self.config = config or {}
        self.ocr_languages = self.config.get('ocr_languages', ['eng'])
        self.max_file_size = self.config.get('max_file_size_mb', 50) * 1024 * 1024  # 50MB default
        
    async def convert_to_markdown(self, file_content: bytes, filename: str, 
                                mime_type: str = None) -> Dict[str, Any]:
        """Convert file content to markdown based on file type"""
        
        if len(file_content) > self.max_file_size:
            return {
                'success': False,
                'error': f'File too large: {len(file_content)} bytes (max: {self.max_file_size})',
                'markdown_content': '',
                'metadata': {}
            }
        
        # Detect MIME type if not provided
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)
            if not mime_type:
                mime_type = 'application/octet-stream'
        
        try:
            # Route to appropriate converter
            if mime_type.startswith('application/pdf'):
                return await self._convert_pdf(file_content, filename)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                             'application/msword']:
                return await self._convert_word(file_content, filename)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                             'application/vnd.ms-powerpoint']:
                return await self._convert_powerpoint(file_content, filename)
            elif mime_type.startswith('image/'):
                return await self._convert_image(file_content, filename)
            elif mime_type in ['text/html', 'application/xhtml+xml']:
                return await self._convert_html(file_content, filename)
            elif mime_type.startswith('text/'):
                return await self._convert_text(file_content, filename)
            else:
                return {
                    'success': False,
                    'error': f'Unsupported file type: {mime_type}',
                    'markdown_content': '',
                    'metadata': {'mime_type': mime_type}
                }
                
        except Exception as e:
            logger.error("Conversion failed", filename=filename, error=str(e))
            return {
                'success': False,
                'error': f'Conversion error: {str(e)}',
                'markdown_content': '',
                'metadata': {'mime_type': mime_type}
            }
    
    async def _convert_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert PDF to markdown"""
        try:
            pdf_document = fitz.open(stream=file_content, filetype="pdf")
            
            markdown_content = f"# {filename}\n\n"
            metadata = {
                'pages': pdf_document.page_count,
                'title': pdf_document.metadata.get('title', ''),
                'author': pdf_document.metadata.get('author', ''),
                'subject': pdf_document.metadata.get('subject', ''),
                'has_images': False,
                'has_tables': False
            }
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    markdown_content += f"## Page {page_num + 1}\n\n{text}\n\n"
                
                # Check for images
                image_list = page.get_images()
                if image_list:
                    metadata['has_images'] = True
                    markdown_content += f"*[Page {page_num + 1} contains {len(image_list)} image(s)]*\n\n"
                
                # Check for tables (basic detection)
                tables = page.find_tables()
                if tables:
                    metadata['has_tables'] = True
                    for table_num, table in enumerate(tables):
                        try:
                            table_data = table.extract()
                            markdown_content += f"### Table {table_num + 1} (Page {page_num + 1})\n\n"
                            
                            # Convert to markdown table
                            if table_data and len(table_data) > 0:
                                # Header row
                                header = table_data[0]
                                markdown_content += "| " + " | ".join(str(cell) if cell else "" for cell in header) + " |\n"
                                markdown_content += "| " + " | ".join("---" for _ in header) + " |\n"
                                
                                # Data rows
                                for row in table_data[1:]:
                                    markdown_content += "| " + " | ".join(str(cell) if cell else "" for cell in row) + " |\n"
                                
                            markdown_content += "\n"
                        except Exception as e:
                            logger.warning("Failed to extract table", page=page_num + 1, table=table_num + 1, error=str(e))
            
            pdf_document.close()
            
            return {
                'success': True,
                'error': None,
                'markdown_content': markdown_content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'PDF conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    async def _convert_word(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert Word document to markdown"""
        try:
            with tempfile.NamedTemporaryFile(suffix='.docx') as temp_file:
                temp_file.write(file_content)
                temp_file.flush()
                
                document = Document(temp_file.name)
                
                markdown_content = f"# {filename}\n\n"
                metadata = {
                    'paragraphs': len(document.paragraphs),
                    'tables': len(document.tables),
                    'has_images': False
                }
                
                # Process paragraphs
                for para in document.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Detect headings based on style
                        if para.style.name.startswith('Heading'):
                            level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                            markdown_content += f"{'#' * (level + 1)} {text}\n\n"
                        else:
                            markdown_content += f"{text}\n\n"
                
                # Process tables
                for table_num, table in enumerate(document.tables):
                    markdown_content += f"## Table {table_num + 1}\n\n"
                    
                    for row_num, row in enumerate(table.rows):
                        cells = [cell.text.strip() for cell in row.cells]
                        
                        if row_num == 0:
                            # Header row
                            markdown_content += "| " + " | ".join(cells) + " |\n"
                            markdown_content += "| " + " | ".join("---" for _ in cells) + " |\n"
                        else:
                            # Data row
                            markdown_content += "| " + " | ".join(cells) + " |\n"
                    
                    markdown_content += "\n"
                
                # Check for images (basic detection)
                for rel in document.part.rels.values():
                    if "image" in rel.target_ref:
                        metadata['has_images'] = True
                        break
                
                return {
                    'success': True,
                    'error': None,
                    'markdown_content': markdown_content,
                    'metadata': metadata
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': f'Word conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    async def _convert_powerpoint(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert PowerPoint presentation to markdown"""
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx') as temp_file:
                temp_file.write(file_content)
                temp_file.flush()
                
                presentation = Presentation(temp_file.name)
                
                markdown_content = f"# {filename}\n\n"
                metadata = {
                    'slides': len(presentation.slides),
                    'has_images': False,
                    'has_tables': False
                }
                
                for slide_num, slide in enumerate(presentation.slides):
                    markdown_content += f"## Slide {slide_num + 1}\n\n"
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # Format as bullet points if multiple lines
                            if '\n' in text:
                                lines = text.split('\n')
                                for line in lines:
                                    if line.strip():
                                        markdown_content += f"- {line.strip()}\n"
                            else:
                                markdown_content += f"{text}\n\n"
                        
                        # Check for tables
                        if shape.has_table:
                            metadata['has_tables'] = True
                            table = shape.table
                            markdown_content += f"### Table (Slide {slide_num + 1})\n\n"
                            
                            for row_num, row in enumerate(table.rows):
                                cells = [cell.text.strip() for cell in row.cells]
                                
                                if row_num == 0:
                                    markdown_content += "| " + " | ".join(cells) + " |\n"
                                    markdown_content += "| " + " | ".join("---" for _ in cells) + " |\n"
                                else:
                                    markdown_content += "| " + " | ".join(cells) + " |\n"
                            
                            markdown_content += "\n"
                    
                    # Check for images
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # Picture type
                            metadata['has_images'] = True
                            markdown_content += f"*[Slide {slide_num + 1} contains image(s)]*\n\n"
                            break
                    
                    markdown_content += "\n"
                
                return {
                    'success': True,
                    'error': None,
                    'markdown_content': markdown_content,
                    'metadata': metadata
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': f'PowerPoint conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    async def _convert_image(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert image to markdown using OCR"""
        try:
            # Open image
            image = Image.open(io.BytesIO(file_content))
            
            metadata = {
                'width': image.width,
                'height': image.height,
                'format': image.format,
                'mode': image.mode
            }
            
            # Preprocess image for better OCR
            processed_image = self._preprocess_image_for_ocr(image)
            
            # Perform OCR
            languages = '+'.join(self.ocr_languages)
            ocr_text = pytesseract.image_to_string(processed_image, lang=languages)
            
            confidence = self._estimate_ocr_confidence(ocr_text)
            metadata['ocr_confidence'] = confidence
            
            markdown_content = f"# {filename}\n\n"
            markdown_content += f"*Image: {image.width}x{image.height} {image.format}*\n\n"
            
            if ocr_text.strip():
                markdown_content += "## Extracted Text\n\n"
                markdown_content += ocr_text.strip()
            else:
                markdown_content += "*No text detected in image*"
            
            return {
                'success': True,
                'error': None,
                'markdown_content': markdown_content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Image conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    async def _convert_html(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert HTML to markdown"""
        try:
            # Detect encoding
            encoding_result = chardet.detect(file_content)
            encoding = encoding_result.get('encoding', 'utf-8')
            
            html_content = file_content.decode(encoding)
            
            # Parse with BeautifulSoup for cleaning
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Convert to markdown
            markdown_content = md(str(soup), heading_style="ATX")
            
            metadata = {
                'encoding': encoding,
                'title': soup.title.string if soup.title else '',
                'links': len(soup.find_all('a')),
                'images': len(soup.find_all('img')),
                'original_length': len(html_content)
            }
            
            return {
                'success': True,
                'error': None,
                'markdown_content': markdown_content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'HTML conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    async def _convert_text(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Convert plain text to markdown (minimal processing)"""
        try:
            # Detect encoding
            encoding_result = chardet.detect(file_content)
            encoding = encoding_result.get('encoding', 'utf-8')
            
            text_content = file_content.decode(encoding)
            
            markdown_content = f"# {filename}\n\n{text_content}"
            
            metadata = {
                'encoding': encoding,
                'lines': len(text_content.split('\n')),
                'chars': len(text_content)
            }
            
            return {
                'success': True,
                'error': None,
                'markdown_content': markdown_content,
                'metadata': metadata
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Text conversion failed: {str(e)}',
                'markdown_content': '',
                'metadata': {}
            }
    
    def _preprocess_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Preprocess image to improve OCR accuracy"""
        try:
            # Convert to grayscale
            if image.mode != 'L':
                image = image.convert('L')
            
            # Convert to numpy array for OpenCV processing
            img_array = np.array(image)
            
            # Apply noise reduction
            img_array = cv2.medianBlur(img_array, 3)
            
            # Apply adaptive thresholding
            img_array = cv2.adaptiveThreshold(
                img_array, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
            )
            
            # Convert back to PIL Image
            processed_image = Image.fromarray(img_array)
            
            # Resize if image is very small
            if processed_image.width < 300 or processed_image.height < 300:
                scale_factor = max(300 / processed_image.width, 300 / processed_image.height)
                new_size = (int(processed_image.width * scale_factor), int(processed_image.height * scale_factor))
                processed_image = processed_image.resize(new_size, Image.Resampling.LANCZOS)
            
            return processed_image
            
        except Exception as e:
            logger.warning("Image preprocessing failed, using original", error=str(e))
            return image
    
    def _estimate_ocr_confidence(self, text: str) -> float:
        """Estimate OCR confidence based on text quality"""
        if not text.strip():
            return 0.0
        
        # Simple heuristics for confidence estimation
        total_chars = len(text)
        if total_chars == 0:
            return 0.0
        
        # Count readable characters
        readable_chars = len(re.findall(r'[a-zA-Z0-9\s.,;:!?]', text))
        readable_ratio = readable_chars / total_chars
        
        # Penalize for too many special characters
        special_chars = len(re.findall(r'[^\w\s.,;:!?-]', text))
        special_ratio = special_chars / total_chars
        
        confidence = readable_ratio - (special_ratio * 0.5)
        return max(0.0, min(1.0, confidence)) 