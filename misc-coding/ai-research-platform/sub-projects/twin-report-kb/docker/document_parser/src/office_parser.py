#!/usr/bin/env python3
"""
Office Parser for Document Parser
Extracts content from Microsoft Office documents (Word, Excel, PowerPoint)
"""

import io
from typing import Dict, List, Any, Optional
import asyncio

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import structlog

logger = structlog.get_logger(__name__)

class OfficeParser:
    """Parses Microsoft Office documents"""
    
    def __init__(self):
        self.supported_formats = ['docx', 'xlsx', 'pptx', 'doc', 'xls', 'ppt']
    
    async def parse(self, content: bytes, content_type: str, options: Dict[str, Any] = None) -> Dict[str, Any]:
        """Parse Office document content"""
        
        if options is None:
            options = {}
        
        logger.info("Starting Office document parsing", 
                   content_type=content_type,
                   content_size=len(content),
                   options=options)
        
        try:
            if content_type == 'docx':
                return await self._parse_word_document(content, options)
            elif content_type == 'xlsx':
                return await self._parse_excel_workbook(content, options)
            elif content_type == 'pptx':
                return await self._parse_powerpoint_presentation(content, options)
            else:
                raise ValueError(f"Unsupported Office format: {content_type}")
                
        except Exception as e:
            logger.error("Office document parsing failed", 
                        content_type=content_type, 
                        error=str(e))
            return {
                "title": f"Office Document Parsing Failed ({content_type})",
                "content": f"Error parsing {content_type} document: {str(e)}",
                "metadata": {"error": str(e), "content_type": content_type},
                "structure": {},
                "stats": {"error": True}
            }
    
    async def _parse_word_document(self, content: bytes, options: Dict[str, Any]) -> Dict[str, Any]:
        """Parse Word document (.docx)"""
        
        try:
            doc_stream = io.BytesIO(content)
            doc = Document(doc_stream)
            
            # Extract text content
            content_parts = []
            paragraphs_info = []
            tables_info = []
            
            # Process paragraphs
            for para_idx, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    content_parts.append(text)
                
                # Analyze paragraph structure
                para_info = {
                    "index": para_idx,
                    "text": text,
                    "style": paragraph.style.name if paragraph.style else None,
                    "alignment": str(paragraph.alignment) if paragraph.alignment else None,
                    "is_heading": self._is_heading_style(paragraph.style.name if paragraph.style else ""),
                    "runs_count": len(paragraph.runs)
                }
                paragraphs_info.append(para_info)
            
            # Process tables
            for table_idx, table in enumerate(doc.tables):
                table_text_parts = []
                table_data = []
                
                for row in table.rows:
                    row_data = []
                    row_text_parts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                        if cell_text:
                            row_text_parts.append(cell_text)
                    
                    if row_text_parts:
                        table_text_parts.append(' | '.join(row_text_parts))
                        table_data.append(row_data)
                
                if table_text_parts:
                    content_parts.append('\n'.join(table_text_parts))
                
                table_info = {
                    "index": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": table_data if options.get("include_table_data", False) else None
                }
                tables_info.append(table_info)
            
            # Extract metadata
            metadata = {
                "parser": "python-docx",
                "content_type": "docx"
            }
            
            # Try to extract document properties
            try:
                core_props = doc.core_properties
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "keywords": core_props.keywords or "",
                    "created": core_props.created.isoformat() if core_props.created else None,
                    "modified": core_props.modified.isoformat() if core_props.modified else None,
                    "last_modified_by": core_props.last_modified_by or "",
                    "revision": core_props.revision or ""
                })
            except Exception as e:
                logger.warning("Failed to extract document properties", error=str(e))
            
            # Combine content
            full_content = '\n\n'.join(content_parts)
            
            result = {
                "title": metadata.get("title") or "Word Document",
                "content": full_content,
                "metadata": metadata,
                "structure": {
                    "paragraphs": paragraphs_info,
                    "tables": tables_info,
                    "sections": len(doc.sections),
                    "styles": [style.name for style in doc.styles]
                },
                "stats": {
                    "word_count": len(full_content.split()),
                    "character_count": len(full_content),
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                    "section_count": len(doc.sections)
                }
            }
            
            logger.info("Word document parsing completed",
                       word_count=result["stats"]["word_count"],
                       paragraph_count=result["stats"]["paragraph_count"])
            
            return result
            
        except Exception as e:
            logger.error("Word document parsing failed", error=str(e))
            raise
    
    async def _parse_excel_workbook(self, content: bytes, options: Dict[str, Any]) -> Dict[str, Any]:
        """Parse Excel workbook (.xlsx)"""
        
        try:
            workbook_stream = io.BytesIO(content)
            workbook = load_workbook(workbook_stream, data_only=True)
            
            content_parts = []
            worksheets_info = []
            
            # Process worksheets
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                
                # Extract data from worksheet
                sheet_data = []
                sheet_text_parts = []
                
                for row in worksheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):
                        sheet_data.append(row_data)
                        sheet_text_parts.append(' | '.join(row_data))
                
                if sheet_text_parts:
                    content_parts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_text_parts))
                
                # Analyze worksheet structure
                sheet_info = {
                    "name": sheet_name,
                    "max_row": worksheet.max_row,
                    "max_column": worksheet.max_column,
                    "data_rows": len(sheet_data),
                    "data": sheet_data if options.get("include_sheet_data", False) else None
                }
                worksheets_info.append(sheet_info)
            
            # Extract metadata
            metadata = {
                "parser": "openpyxl",
                "content_type": "xlsx",
                "worksheet_count": len(workbook.sheetnames),
                "worksheet_names": workbook.sheetnames
            }
            
            # Try to extract workbook properties
            try:
                props = workbook.properties
                metadata.update({
                    "title": props.title or "",
                    "creator": props.creator or "",
                    "description": props.description or "",
                    "subject": props.subject or "",
                    "keywords": props.keywords or "",
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                    "last_modified_by": props.lastModifiedBy or ""
                })
            except Exception as e:
                logger.warning("Failed to extract workbook properties", error=str(e))
            
            # Combine content
            full_content = '\n\n'.join(content_parts)
            
            result = {
                "title": metadata.get("title") or "Excel Workbook",
                "content": full_content,
                "metadata": metadata,
                "structure": {
                    "worksheets": worksheets_info
                },
                "stats": {
                    "word_count": len(full_content.split()),
                    "character_count": len(full_content),
                    "worksheet_count": len(workbook.sheetnames),
                    "total_cells": sum(info["data_rows"] * info["max_column"] for info in worksheets_info)
                }
            }
            
            logger.info("Excel workbook parsing completed",
                       worksheet_count=result["stats"]["worksheet_count"],
                       total_cells=result["stats"]["total_cells"])
            
            return result
            
        except Exception as e:
            logger.error("Excel workbook parsing failed", error=str(e))
            raise
    
    async def _parse_powerpoint_presentation(self, content: bytes, options: Dict[str, Any]) -> Dict[str, Any]:
        """Parse PowerPoint presentation (.pptx)"""
        
        try:
            ppt_stream = io.BytesIO(content)
            presentation = Presentation(ppt_stream)
            
            content_parts = []
            slides_info = []
            
            # Process slides
            for slide_idx, slide in enumerate(presentation.slides):
                slide_text_parts = []
                shapes_info = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                    
                    # Analyze shape
                    shape_info = {
                        "type": str(shape.shape_type),
                        "has_text": hasattr(shape, "text") and bool(shape.text.strip()),
                        "text_length": len(shape.text) if hasattr(shape, "text") else 0
                    }
                    shapes_info.append(shape_info)
                
                if slide_text_parts:
                    slide_content = '\n'.join(slide_text_parts)
                    content_parts.append(f"Slide {slide_idx + 1}:\n{slide_content}")
                
                # Analyze slide structure
                slide_info = {
                    "index": slide_idx + 1,
                    "shapes_count": len(slide.shapes),
                    "text_shapes": len([s for s in shapes_info if s["has_text"]]),
                    "layout": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
                    "shapes": shapes_info if options.get("include_shape_details", False) else None
                }
                slides_info.append(slide_info)
            
            # Extract metadata
            metadata = {
                "parser": "python-pptx",
                "content_type": "pptx",
                "slide_count": len(presentation.slides)
            }
            
            # Try to extract presentation properties
            try:
                core_props = presentation.core_properties
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "keywords": core_props.keywords or "",
                    "created": core_props.created.isoformat() if core_props.created else None,
                    "modified": core_props.modified.isoformat() if core_props.modified else None,
                    "last_modified_by": core_props.last_modified_by or "",
                    "revision": core_props.revision or ""
                })
            except Exception as e:
                logger.warning("Failed to extract presentation properties", error=str(e))
            
            # Combine content
            full_content = '\n\n'.join(content_parts)
            
            result = {
                "title": metadata.get("title") or "PowerPoint Presentation",
                "content": full_content,
                "metadata": metadata,
                "structure": {
                    "slides": slides_info,
                    "slide_layouts": [layout.name for layout in presentation.slide_layouts]
                },
                "stats": {
                    "word_count": len(full_content.split()),
                    "character_count": len(full_content),
                    "slide_count": len(presentation.slides),
                    "total_shapes": sum(info["shapes_count"] for info in slides_info),
                    "text_shapes": sum(info["text_shapes"] for info in slides_info)
                }
            }
            
            logger.info("PowerPoint presentation parsing completed",
                       slide_count=result["stats"]["slide_count"],
                       word_count=result["stats"]["word_count"])
            
            return result
            
        except Exception as e:
            logger.error("PowerPoint presentation parsing failed", error=str(e))
            raise
    
    def _is_heading_style(self, style_name: str) -> bool:
        """Check if a style is a heading style"""
        if not style_name:
            return False
        
        heading_indicators = ['heading', 'title', 'subtitle', 'header']
        style_lower = style_name.lower()
        
        return any(indicator in style_lower for indicator in heading_indicators)
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported formats"""
        return self.supported_formats 